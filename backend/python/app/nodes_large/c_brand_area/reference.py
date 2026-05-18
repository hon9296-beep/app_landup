"""
브랜드 메뉴얼 규칙 추출 노드 — Shin 코드 베이스.

브랜드 메뉴얼 PDF → 규칙 추출 (Agent 1).
레퍼런스 이미지 로드는 ref_image_loader.py로 분리됨.
"""
import logging
import os
import re

import fitz
from anthropic import Anthropic

from app.state import LargeState
from app.utils import normalize_object_type, OBJECT_STANDARDS
from app.vmd_constants import (
    VMD_BOUNDARIES,
    VMD_PAIR_RULES,
)
# 2026-05-03 — large 자유 디자인 정합성 위해 MAX_COUNT_BY_CATEGORY / scale_count fallback 폐기.
# brand LLM 이 max_count 누락 시 design LLM 자율 결정 (verify 노드가 안전망).
# fallback 전략 재구상은 디자인 참조 로직 + 배치 알고리즘 손볼 때 (TR_S 고도화 트랙).
# small 은 vmd_constants 의 정의 그대로 사용 (rendy 영역, rule 18).

logger = logging.getLogger(__name__)

# 브랜드 기본값
BRAND_DEFAULTS = {
    "clearspace_mm": 1000, "logo_clearspace_mm": 500,
    "character_orientation": "자유", "prohibited_material": None,
}
FIRE_RULES = {"main_corridor_min_mm": 900, "emergency_path_min_mm": 1200}
CONSTRUCTION_RULES = {"wall_clearance_mm": 300, "object_gap_mm": 300}

# ── VMD 상수는 app.vmd_constants로 통합 (2026-04-19, 결정로그 ①) ──
# VMD_BOUNDARIES, VMD_PAIR_RULES, VMD_WALL_ATTACHMENT,
# MAX_COUNT_CHARACTER_IP, MAX_COUNT_BY_CATEGORY, scale_count
# 모두 app.vmd_constants에서 import.
# 로컬 정의 폐기 — small과 동일한 정본 사용.


def run(state: LargeState) -> LargeState:
    """브랜드 메뉴얼 → 규칙 추출 (이미지 로드 없음)."""
    brand_bytes = state.get("brand_bytes")
    file_type = state.get("brand_file_type", "pdf")
    if brand_bytes:
        brand_data = _run_brand_agent(brand_bytes, file_type=file_type)
    else:
        brand_data = _fallback_brand_defaults()

    return {"brand_data": brand_data}


# ── 브랜드 추출 ───────────────────────────────────────────────────────────

BRAND_SYSTEM = """당신은 팝업스토어 브랜드 메뉴얼에서 공간 배치 제약 조건을 추출하는 전문가입니다.
정확히 명시된 내용만 추출하고, 문서에 없는 내용은 절대 추측하지 마세요.

## 단위 변환 규칙 (항상 적용)
- cm → mm: ×10 (예: 50cm → 500)
- m  → mm: ×1000 (예: 1.2m → 1200)
- clearspace_mm, logo_clearspace_mm가 범위로 표현된 경우 (예: 300~500mm) → 최솟값 사용 (예: 300)

## confidence 판단 기준
- high:   수치나 명칭이 문서에 명확히 적혀있음
- medium: 언급은 있지만 해석이 필요한 경우
- low:    암시적이거나 불분명한 경우, 또는 문서에 없어서 null인 경우

## 수량 규칙 (min_count, max_count)
- 문서에 숫자로 명확히 적혀있을 때만 설정 (예: "최소 2개", "3개 이상")
- 명시되지 않은 경우 반드시 null — 절대 추측하지 말 것
- "메인", "대표", "주요" 같은 수식어는 수량 근거가 아님 → null

## 오브젝트 타입 정의 (매뉴얼 항목 분류 기준)
- character_bbox: 바닥에 독립으로 서는 3D 입체물. 브랜드 마스코트·캐릭터·피규어·오브제·설치물 등 시각적 포인트가 되는 모든 입체 오브젝트
- photo_wall: 배경 패널·포토월·가벽 등 사진 촬영 배경으로 쓰이는 평면/벽면 구조물
- counter: 직원이 서서 고객 응대·결제를 처리하는 테이블형 구조물
- display_table: 바닥에 독립으로 서는 상품 진열 테이블 (아일랜드형)
- shelf_wall: 벽에 부착하거나 벽에 세우는 상품 진열 선반
- shelf_3tier: 바닥에 독립으로 서는 다단 선반
- banner_stand: 천·패브릭·인쇄물을 걸어두는 세로형 스탠드
- signage_stand: A형 입간판, 방향·정보 안내 표지판
- kiosk: 무인 결제기·터치스크린 안내 단말기
- partition_wall_I: 공간 분리용 일자형 임시 벽체
- partition_wall_L: 공간 분리용 ㄱ자형 코너 벽체"""

BRAND_PROMPT = """아래 브랜드 매뉴얼에서 배치 규칙을 추출하세요.
문서에 명시되지 않은 항목은 value를 null, confidence를 low로 설정하세요.
수치는 반드시 mm 정수로 변환하세요."""

BRAND_TOOL = {
    "name": "extract_brand_rules",
    "description": "브랜드 매뉴얼에서 공간 배치 제약 조건을 구조화하여 추출합니다.",
    "input_schema": {
        "type": "object",
        "properties": {
            "brand_category": {
                "type": "string",
                "enum": ["캐릭터 IP", "패션 브랜드", "F&B", "뷰티·코스메틱", "기타"],
                "description": "브랜드 카테고리. 명시되지 않으면 '기타'.",
            },
            "clearspace_mm": {
                "type": "object",
                "properties": {
                    "value": {"type": ["integer", "null"], "description": "이격 거리 (mm 정수). cm→×10, m→×1000 변환."},
                    "confidence": {"type": "string", "enum": ["low", "medium", "high"]},
                },
                "required": ["value", "confidence"],
            },
            "character_orientation": {
                "type": "object",
                "properties": {
                    "value": {"type": ["string", "null"], "enum": ["입구 정면", "벽면", "자유", None]},
                    "confidence": {"type": "string", "enum": ["low", "medium", "high"]},
                },
                "required": ["value", "confidence"],
            },
            "prohibited_material": {
                "type": "object",
                "properties": {
                    "value": {"type": ["string", "null"], "description": "금지 소재. 없으면 null."},
                    "confidence": {"type": "string", "enum": ["low", "medium", "high"]},
                },
                "required": ["value", "confidence"],
            },
            "logo_clearspace_mm": {
                "type": "object",
                "properties": {
                    "value": {"type": ["integer", "null"], "description": "로고 여백 (mm 정수). 없으면 null."},
                    "confidence": {"type": "string", "enum": ["low", "medium", "high"]},
                },
                "required": ["value", "confidence"],
            },
            "pair_rules": {
                "type": "array",
                "description": "오브젝트 쌍 간 관계 규정. 문서에 명시된 것만. 없으면 빈 배열.",
                "items": {
                    "type": "object",
                    "properties": {
                        "object_a": {
                            "type": "string",
                            "enum": [
                                "counter", "display_table", "character_bbox",
                                "photo_wall", "shelf_wall", "shelf_3tier",
                                "banner_stand", "signage_stand", "kiosk",
                                "partition_wall_I", "partition_wall_L",
                            ],
                            "description": "표준 오브젝트 타입 ID로 매핑.",
                        },
                        "object_b": {
                            "type": "string",
                            "enum": [
                                "counter", "display_table", "character_bbox",
                                "photo_wall", "shelf_wall", "shelf_3tier",
                                "banner_stand", "signage_stand", "kiosk",
                                "partition_wall_I", "partition_wall_L",
                            ],
                            "description": "표준 오브젝트 타입 ID로 매핑.",
                        },
                        "relation": {"type": "string", "enum": ["join", "separate", "adjacent"]},
                        "min_gap_mm": {"type": "integer", "description": "최소 간격 (join이면 0)"},
                    },
                    "required": ["object_a", "object_b", "relation", "min_gap_mm"],
                },
            },
            "figures_mentioned": {
                "type": "array",
                "description": "매뉴얼에 언급된 독립 입체 조형물 목록. 상품 판매·진열 목적이 아닌, 시각적 포인트·포토존·공간 연출 목적으로 공간에 단독 배치되는 3D 오브젝트. 캐릭터 조형물·마스코트 피규어·아트 오브제·브랜드 상징물 등이 해당. 진열대·선반·계산대·배너·가벽은 제외. 없으면 빈 배열.",
                "items": {"type": "string"},
            },
            "concept_areas_hint": {
                "type": "array",
                "description": (
                    "매뉴얼에 명시된 영역(zone) 목록. '맞이존' / '스킨케어 코너' / '체험 라운지' 같은 "
                    "공간 영역 의미. 매뉴얼에 zone / 영역 / 코너 / 존 같은 키워드 + 위치/역할 명시되어 있을 때만 추출. "
                    "없으면 빈 배열. 기본 8종 (맞이/포토/체험/상영/굿즈판매/결제/혼합/휴식) 외 커스텀 영역도 OK. "
                    "2026-05-06 burning_task 2단계 — 세 번째 작업 (영역 dynamic) 신설."
                ),
                "items": {
                    "type": "object",
                    "properties": {
                        "name_ko": {
                            "type": "string",
                            "description": "한국어 영역명 (매뉴얼에 명시된 그대로 — 예: '스킨케어 코너', '캐릭터 맞이존')",
                        },
                        "name_en": {
                            "type": "string",
                            "description": (
                                "영문 키 (snake_case, DB 저장용). LLM 가 ko 보고 자동 변환. "
                                "기본 8종 (welcome/photo/experience/screening/retail/checkout/hybrid/lounge) 시 그대로, "
                                "커스텀 영역은 의미 매칭해서 새 영문 키 (예: '스킨케어 코너' → 'skincare_corner')."
                            ),
                        },
                        "description": {
                            "type": "string",
                            "description": "영역 역할 / 디자인 톤 (1~2줄, 매뉴얼 명시 또는 LLM 추론)",
                        },
                        "target_objects": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": (
                                "이 영역에 권장되는 오브젝트 종류 (표준 ID list). 매뉴얼 명시 시 그대로 매핑, "
                                "명시 X 시 빈 배열 (코드가 가장 비슷한 8종 영역의 target_objects fallback)."
                            ),
                        },
                    },
                    "required": ["name_ko", "name_en", "description"],
                },
            },
            "placement_rules": {
                "type": "array",
                "description": "집기별 배치 규정. 문서에 명시된 것만. 없으면 빈 배열.",
                "items": {
                    "type": "object",
                    "properties": {
                        "object_type": {
                            "type": "string",
                            "enum": [
                                "counter", "display_table", "character_bbox",
                                "photo_wall", "shelf_wall", "shelf_3tier",
                                "banner_stand", "signage_stand", "kiosk",
                                "partition_wall_I", "partition_wall_L",
                            ],
                            "description": "표준 오브젝트 타입 ID. 매뉴얼의 명칭을 가장 가까운 표준 ID로 매핑.",
                        },
                        "name": {"type": "string"},
                        "preferred_wall": {"type": ["string", "null"]},
                        "required_direction": {"type": ["string", "null"]},
                        "width_mm": {"type": ["integer", "null"]},
                        "depth_mm": {"type": ["integer", "null"]},
                        "height_mm": {"type": ["integer", "null"]},
                        "min_count": {"type": ["integer", "null"]},
                        "max_count": {"type": ["integer", "null"]},
                        "max_count_source": {
                            "type": "string",
                            "enum": ["manual", "inferred"],
                            "description": "manual: 매뉴얼에 숫자로 명확히 명시됨. inferred: 명시 없이 추측한 값.",
                        },
                    },
                    "required": ["object_type"],
                },
            },
        },
        "required": [
            "brand_category", "clearspace_mm", "character_orientation",
            "prohibited_material", "logo_clearspace_mm", "pair_rules", "placement_rules",
            "figures_mentioned",
            # 2026-05-06 burning_task 2단계 세 번째 작업 — 영역 dynamic 신설.
            "concept_areas_hint",
        ],
    },
}


def _clamp_dimensions(rule: dict) -> None:
    """LLM 추출 규격을 VMD_BOUNDARIES 기준으로 교정 (in-place).

    - w/d/h 중 하나라도 누락 → std 세트 전체로 덮어쓰기
    - 값이 있어도 min 미만 또는 max 초과 → std로 강제 교정
    - max_count fallback 은 2026-05-03 폐기 (large 자유 디자인 정합)
    """
    obj_type = rule.get("object_type", "")
    bounds = VMD_BOUNDARIES.get(obj_type)

    # ── w/d/h 교정 ──
    if bounds:
        dims = ("width_mm", "depth_mm", "height_mm")

        # 하나라도 누락이면 세트 통째로 std 적용
        if not all(rule.get(d) for d in dims):
            for d in dims:
                rule[d] = bounds[d]["std"]
            logger.info(f"[reference] {obj_type}: 규격 누락 → std 세트 적용 "
                        f"({bounds['width_mm']['std']}×{bounds['depth_mm']['std']}×{bounds['height_mm']['std']})")
        else:
            # 개별 필드 범위 검사 — 벗어나면 std로 교정
            for d in dims:
                val = rule[d]
                b = bounds[d]
                if val < b["min"] or val > b["max"]:
                    logger.info(f"[reference] {obj_type}.{d}: {val} → {b['std']} "
                                f"(범위 {b['min']}~{b['max']} 벗어남)")
                    rule[d] = b["std"]

    # ── max_count 교정: 2026-05-03 폐기 ──
    # 옛 fallback (`MAX_COUNT_BY_CATEGORY` + `_scale_count`) 은 large 자유 디자인 정합 X.
    # brand LLM 이 max_count 누락 시 design LLM 자율 결정 (verify 노드가 안전망).
    # 누락 시 object_selection.py:52 가 `count = ... or 1` 로 default 1 박음.
    # fallback 전략 재구상은 디자인 참조 로직 + 배치 알고리즘 손볼 때.


def _run_brand_agent(file_bytes: bytes, file_type: str = "pdf") -> dict:
    if file_type == "pptx":
        text = _extract_pptx_text(file_bytes)
    elif file_type == "docx":
        text = _extract_docx_text(file_bytes)
    elif file_type == "xlsx":
        text = _extract_xlsx_text(file_bytes)
    else:
        text = _extract_pdf_text(file_bytes)

    if not text.strip():
        return _fallback_brand_defaults()

    # Regex 주석 삽입 — LLM 단위 변환 실수 방지
    text = _annotate_measurements(text)

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        logger.info("[reference] API 키 없음 — 기본값 반환")
        return _fallback_brand_defaults()

    client = Anthropic(api_key=api_key)

    for attempt in range(3):
        try:
            response = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=3072,
                system=BRAND_SYSTEM,
                tools=[BRAND_TOOL],
                tool_choice={"type": "tool", "name": "extract_brand_rules"},
                messages=[{"role": "user", "content": f"{BRAND_PROMPT}\n\n## 메뉴얼\n\n{text}"}],
            )
            from app.token_tracker import track_usage
            track_usage("large.reference", response)
            if not response.content:
                continue
            tool_block = next((b for b in response.content if b.type == "tool_use"), None)
            if tool_block is None:
                continue
            raw = tool_block.input

            # defaults merge
            for key, default in BRAND_DEFAULTS.items():
                if key in raw and isinstance(raw[key], dict):
                    if raw[key].get("value") is None and default is not None:
                        raw[key]["value"] = default
                        raw[key]["source"] = "default"

            # placement_rules VMD_BOUNDARIES 기반 교정
            for rule in raw.get("placement_rules", []):
                _clamp_dimensions(rule)

            # pair_rules: LLM 추출분 + VMD 기본 규칙 병합
            extracted_pairs = raw.get("pair_rules") or raw.get("relationships") or []
            merged_pairs = _merge_pair_rules(extracted_pairs)

            brand_fields = {k: v for k, v in raw.items()
                            if k not in ("placement_rules", "pair_rules", "relationships")}

            normalized_rules = _normalize_placement_rules(raw.get("placement_rules", []))

            # character_bbox max_count 보정 (정규화 이후: 중복 합산 방지)
            # 1순위: max_count_source == "manual" → 매뉴얼 명시값 그대로 사용
            # 2순위: max_count_source == "inferred" + figures_mentioned 있음 → figures 개수로 교체
            # 3순위: max_count == null → 2026-05-03 fallback 폐기 (design LLM 자율 + verify 안전망)
            figures = raw.get("figures_mentioned", [])
            for rule in normalized_rules:
                if rule.get("object_type") == "character_bbox":
                    source = rule.get("max_count_source", "inferred")
                    if source == "inferred" and figures:
                        rule["max_count"] = len(figures)
                        logger.info(f"[reference] character_bbox max_count → {len(figures)} (figures_mentioned 기준)")

            return {
                "brand": brand_fields,
                "fire": FIRE_RULES,
                "construction": CONSTRUCTION_RULES,
                "placement_rules": normalized_rules,
                "pair_rules": merged_pairs,
            }
        except Exception as e:
            logger.warning(f"[reference] brand attempt {attempt+1} 실패: {e}")

    return _fallback_brand_defaults()


def _merge_pair_rules(extracted: list) -> list:
    """LLM 추출 pair_rules + VMD 기본 규칙 병합.

    LLM 추출분이 우선. 같은 (object_a, object_b) 쌍이 VMD에도 있으면 LLM 것을 유지.
    """
    # 추출분 구조화 (자연어 형식이면 변환)
    structured = []
    for rule in extracted:
        if isinstance(rule, dict) and "object_a" in rule and "relation" in rule:
            structured.append(rule)
        elif isinstance(rule, dict) and "rule" in rule:
            # Rendy 호환: {"rule": "자연어"} → 구조화 불가, 스킵
            logger.info(f"[reference] pair_rule 자연어 형식 스킵: {rule.get('rule', '')[:50]}")
        elif isinstance(rule, str):
            logger.info(f"[reference] pair_rule 문자열 형식 스킵: {rule[:50]}")

    # 추출분에 있는 쌍 키 수집
    extracted_pairs = set()
    for r in structured:
        extracted_pairs.add((r["object_a"], r["object_b"]))
        extracted_pairs.add((r["object_b"], r["object_a"]))  # 양방향

    # VMD 기본 규칙 중 추출분과 겹치지 않는 것만 추가
    for vmd_rule in VMD_PAIR_RULES:
        key = (vmd_rule["object_a"], vmd_rule["object_b"])
        if key not in extracted_pairs:
            structured.append({**vmd_rule, "source": "vmd_default"})

    logger.info(f"[reference] pair_rules: {len(extracted)}개 추출 → {len(structured)}개 (VMD 병합)")
    return structured


def _fallback_brand_defaults() -> dict:
    brand = {}
    for k, v in BRAND_DEFAULTS.items():
        brand[k] = {"value": v, "confidence": None, "source": "default"}
    brand["brand_category"] = "\uae30\ud0c0"
    return {
        "brand": brand,
        "fire": FIRE_RULES,
        "construction": CONSTRUCTION_RULES,
        "placement_rules": [],
        "pair_rules": [{**r, "source": "vmd_default"} for r in VMD_PAIR_RULES],
    }


def _normalize_placement_rules(rules: list) -> list:
    """LLM이 추출한 placement_rules의 object_type을 표준 ID로 정규화.

    메뉴얼에 "포토존 배경 월"이라고 적혀있어도 → "photo_wall"으로 통일.
    치수가 없으면 OBJECT_STANDARDS에서 기본값 채움.
    """
    normalized = []
    seen_types = set()
    for rule in rules:
        raw_type = rule.get("object_type", "")
        std_id = normalize_object_type(raw_type)
        std = OBJECT_STANDARDS.get(std_id)

        # 표준 치수 보충 — VMD_BOUNDARIES 우선, 없으면 OBJECT_STANDARDS fallback
        rule["object_type"] = std_id
        bounds = VMD_BOUNDARIES.get(std_id)
        if bounds:
            rule.setdefault("width_mm", bounds["width_mm"]["std"])
            rule.setdefault("depth_mm", bounds["depth_mm"]["std"])
            rule.setdefault("height_mm", bounds["height_mm"]["std"])
        if std:
            rule.setdefault("name", std["name"])

        # 같은 표준 ID 중복 방지 (max_count로 수량 관리)
        if std_id in seen_types:
            # 기존 것의 max_count 증가
            for existing in normalized:
                if existing["object_type"] == std_id:
                    existing["max_count"] = existing.get("max_count", 1) + rule.get("max_count", 1)
                    break
        else:
            seen_types.add(std_id)
            normalized.append(rule)

    logger.info(f"[reference] placement_rules 정규화: {len(rules)}개 → {len(normalized)}개 표준 타입")
    return normalized


def _annotate_measurements(text: str) -> str:
    """PDF 텍스트에 단위 변환 주석 삽입 (Regex 단계).

    LLM이 단위 변환 실수를 줄이도록 원본 표현 옆에 mm 변환값을 추가.
    순서 중요: 범위 패턴 먼저, 단일 패턴 나중에 처리.

    예: "50cm"      → "50cm[=500mm]"
        "1.2m"      → "1.2m[=1200mm]"
        "30~50cm"   → "30~50cm[=300~500mm, 최솟값=300mm]"
        "300~500mm" → "300~500mm[최솟값=300mm]"
    """
    # 1. cm 범위: "30~50cm" → "30~50cm[=300~500mm, 최솟값=300mm]"
    def replace_cm_range(m):
        lo_mm = int(float(m.group(1)) * 10)
        hi_mm = int(float(m.group(2)) * 10)
        return f"{m.group(0)}[={lo_mm}~{hi_mm}mm, 최솟값={lo_mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*[~～]\s*(\d+(?:\.\d+)?)\s*cm(?!\[)', replace_cm_range, text)

    # 2. m 범위 (mm 제외): "1~2m" → "1~2m[=1000~2000mm, 최솟값=1000mm]"
    def replace_m_range(m):
        lo_mm = int(float(m.group(1)) * 1000)
        hi_mm = int(float(m.group(2)) * 1000)
        return f"{m.group(0)}[={lo_mm}~{hi_mm}mm, 최솟값={lo_mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*[~～]\s*(\d+(?:\.\d+)?)\s*m(?!m|\[|²)', replace_m_range, text)

    # 3. mm 범위: "300~500mm" → "300~500mm[최솟값=300mm]"
    def replace_mm_range(m):
        lo_mm = int(m.group(1))
        return f"{m.group(0)}[최솟값={lo_mm}mm]"
    text = re.sub(r'(\d+)\s*[~～]\s*(\d+)\s*mm(?!\[)', replace_mm_range, text)

    # 4. 단일 cm: "50cm" → "50cm[=500mm]"
    def replace_cm(m):
        mm = int(float(m.group(1)) * 10)
        return f"{m.group(0)}[={mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*cm(?!\[)', replace_cm, text)

    # 5. 단일 m (mm·m²·주석 제외): "1.2m" → "1.2m[=1200mm]"
    def replace_m(m):
        mm = int(float(m.group(1)) * 1000)
        return f"{m.group(0)}[={mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*m(?!m|\[|²)', replace_m, text)

    return text


def _extract_pdf_text(pdf_bytes: bytes) -> str:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page_texts = []
    for page in doc:
        texts = []
        table_bboxes = []

        try:
            for table in page.find_tables():
                table_bboxes.append(table.bbox)
                rows = []
                for row in table.extract():
                    cells = [str(c).strip() for c in row if c and str(c).strip()]
                    if cells:
                        rows.append(" ".join(cells))
                if rows:
                    texts.append("\n".join(rows))
        except Exception:
            pass

        for block in page.get_text("blocks", sort=True):
            if block[6] != 0:
                continue
            text = block[4].strip()
            if not text:
                continue
            bx0, by0, bx1, by1 = block[:4]
            in_table = any(
                bx0 < tx1 and bx1 > tx0 and by0 < ty1 and by1 > ty0
                for (tx0, ty0, tx1, ty1) in table_bboxes
            )
            if not in_table:
                texts.append(text)

        if texts:
            page_texts.append("\n".join(texts))

    doc.close()
    return "\n\n".join(page_texts)


def _extract_pptx_text(data: bytes) -> str:
    import io
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _extract_shapes(shapes, texts):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _extract_shapes(shape.shapes, texts)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line:
                        texts.append(line)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" ".join(cells))

    prs = Presentation(io.BytesIO(data))
    texts = []
    for slide in prs.slides:
        _extract_shapes(slide.shapes, texts)
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            note = notes_tf.text.strip()
            if note:
                texts.append(note)
        except Exception:
            pass
    return "\n".join(texts)


def _extract_docx_text(data: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(data))
    texts = []
    for section in doc.sections:
        for hf in (section.header, section.footer):
            for para in hf.paragraphs:
                if para.text.strip():
                    texts.append(para.text.strip())
    texts += [para.text for para in doc.paragraphs if para.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                texts.append(" ".join(cells))
    return "\n".join(texts)


def _extract_xlsx_text(data: bytes) -> str:
    import io
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    texts = []
    for ws in wb.worksheets:
        texts.append(f"[{ws.title}]")
        for row in ws.iter_rows(values_only=True):
            line = " ".join(str(cell) for cell in row if cell is not None and str(cell).strip())
            if line:
                texts.append(line)
    wb.close()
    return "\n".join(texts)
