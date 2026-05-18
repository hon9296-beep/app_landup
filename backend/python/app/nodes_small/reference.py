"""
브랜드 메뉴얼 규칙 추출 노드 — Shin 코드 베이스.

브랜드 메뉴얼 PDF → 규칙 추출 (Agent 1).
레퍼런스 이미지 로드는 ref_image_loader.py로 분리됨.
"""
import logging
import os
import re
from typing import Optional

import fitz
from anthropic import Anthropic
from pydantic import Field

from app.nodes_small.llm_policy import StrictLLMModel
from app.state import SmallState
from app.utils import parse_llm_json, normalize_object_type, OBJECT_STANDARDS
from app.vmd_constants import (
    VMD_PAIR_RULES,
    MAX_COUNT_CHARACTER_IP as _MAX_COUNT_CHARACTER_IP,
    MAX_COUNT_BY_CATEGORY as _MAX_COUNT_BY_CATEGORY,
    scale_count as _scale_count,
    get_vmd_boundaries,
)

logger = logging.getLogger(__name__)


# 2026-04-29 Phase 2 harness: BRAND_TOOL.input_schema 의 top-level 필드와 1:1 대응.
# nested dict (clearspace_mm 등 confidence 포함) / list 는 dict/list 그대로 받음 — 후처리에서
# value/confidence 키 접근. StrictLLMModel: extra="allow" + 위험 키 거부 + 새 필드 로깅.
#
# 2026-05-01 (G 옵션) — model_validator 제거.
# Why (어제 4-30 c6f3dd9 의 잘못된 결합 fix):
#   PR #366 의 model_validator 가 placement_rules 핵심 풀 누락 시 ValueError raise →
#   llm_harness 가 LLMSchemaValidationError 로 wrap → 3회 retry 모두 실패 → 기존 except
#   절이 LLMHarnessError 까지 잡아서 _fallback_brand_defaults() 호출 → brand 전체 default
#   리셋. placement_rules 한 부분 검증 실패가 brand_category / clearspace_mm 등 LLM 정상
#   추출 정보까지 모두 폐기시키는 잘못된 결합. fallback 함수의 역할 (LLM 호출 자체 실패
#   안전망) 을 넘어 검증 실패 시까지 호출되는 것은 권한 밖.
# After: model_validator 제거. placement_rules 부실 검사는 _run_brand_agent 후처리에서
#   logger.warning 으로만 수행 (raise 안 함). brand 응답 LLM 추출분 항상 보존.
class BrandRulesResult(StrictLLMModel):
    brand_category: str = "기타"
    clearspace_mm: Optional[dict] = None
    character_orientation: Optional[dict] = None
    prohibited_material: Optional[dict] = None
    logo_clearspace_mm: Optional[dict] = None
    pair_rules: list = Field(default_factory=list)
    relationships: list = Field(default_factory=list)  # legacy alias — _merge_pair_rules 가 fallback 으로 봄
    figures_mentioned: list = Field(default_factory=list)
    placement_rules: list = Field(default_factory=list)

# 브랜드 기본값
BRAND_DEFAULTS = {
    "clearspace_mm": 600, "logo_clearspace_mm": 500,  # 인간 실질 활동 반경 600mm
    "character_orientation": "자유", "prohibited_material": None,
}
FIRE_RULES = {"main_corridor_min_mm": 900, "emergency_path_min_mm": 1200}
CONSTRUCTION_RULES = {"wall_clearance_mm": 300, "object_gap_mm": 300}

# ── VMD 상수는 constants.py로 이동 (2026-04-16) ──────────────────────────
# VMD_BOUNDARIES, VMD_PAIR_RULES, VMD_WALL_ATTACHMENT,
# MAX_COUNT_CHARACTER_IP, MAX_COUNT_BY_CATEGORY, scale_count
# → from app.vmd_constants import ... 으로 사용


def run(state: SmallState) -> SmallState:
    """브랜드 메뉴얼 → 규칙 추출 (이미지 로드 없음)."""
    brand_bytes = state.get("brand_bytes")
    file_type = state.get("brand_file_type", "pdf")
    if brand_bytes:
        brand_data = _run_brand_agent(brand_bytes, file_type=file_type)
    else:
        brand_data = _fallback_brand_defaults()

    return {"brand_data": brand_data}


# ── 브랜드 추출 ───────────────────────────────────────────────────────────

# BRAND prompt / Tool schema — #491 prompts 중앙화 (nodes_small/prompts/reference.py)
from app.nodes_small.prompts.reference import (
    BRAND_SYSTEM,
    BRAND_PROMPT,
    BRAND_TOOL,
)


def _log_minimal_placement_rules_warning(brand_category: str, placement_rules: list) -> None:
    """G 옵션 후처리 모니터링 — minimal 풀 누락 시 logger.warning.

    PR #366 의 model_validator (raise → fallback) 잘못된 결합 fix. raise 안 함.
    부실 응답이어도 brand 응답 LLM 추출분은 보존하고 다음 단계로 진행.

    Args:
        brand_category: LLM 이 추출한 brand_category 값
        placement_rules: LLM 응답의 placement_rules 리스트 (raw)

    Behavior:
        - SSOT (app.categories) 의 카테고리별 minimal 풀 lookup
        - 빈 set 이면 모니터링 skip (테크/아트/엔터 등 미등록 카테고리)
        - raw object_type 을 normalize_object_type 으로 표준 ID 매핑 후 set 비교
        - 누락 발견 시 logger.warning. raise 절대 X.
    """
    from app.categories import get_category

    required = get_category(brand_category).minimal_placement_rules
    if not required:
        return  # SSOT 미등록 카테고리 — 모니터링 skip

    actual: set[str] = {
        normalize_object_type(r.get("object_type", ""))
        for r in placement_rules
        if isinstance(r, dict) and r.get("object_type")
    }
    missing = required - actual
    if missing:
        logger.warning(
            f"[reference] placement_rules 핵심 풀 누락 (모니터링) — "
            f"brand_category='{brand_category}', 기대={sorted(required)}, "
            f"실제={sorted(actual)}, 누락={sorted(missing)}. "
            f"매뉴얼에 해당 기물 정보 부재 추정 — brand 응답은 그대로 보존."
        )


def _clamp_dimensions(rule: dict, brand_category: str = "기타", usable_area_m2: float = 50.0) -> None:
    """LLM 추출 규격 + max_count를 VMD_BOUNDARIES 기준으로 교정 (in-place).

    - w/d/h 중 하나라도 누락 → std 세트 전체로 덮어쓰기
    - 값이 있어도 min 미만 또는 max 초과 → std로 강제 교정
    - max_count 누락 시 카테고리별 기본값 × 면적 스케일링 적용
    """
    obj_type = rule.get("object_type", "")
    bounds = get_vmd_boundaries(brand_category).get(obj_type)

    # ── w/d/h 교정 ──
    if bounds:
        dims = ("width_mm", "depth_mm", "height_mm")

        # 하나라도 누락이면 세트 통째로 std 적용
        if not all(rule.get(d) for d in dims):
            for d in dims:
                rule[d] = bounds[d]["std"]
            logger.info(f"[reference] {obj_type}: 규격 누락 → std 세트 적용 "
                        f"({bounds['width_mm']['std']}×{bounds['depth_mm']['std']}×{bounds['height_mm']['std']})")
        else:
            # 개별 필드 범위 검사 — 벗어나면 std로 교정
            for d in dims:
                val = rule[d]
                b = bounds[d]
                if val < b["min"] or val > b["max"]:
                    logger.info(f"[reference] {obj_type}.{d}: {val} → {b['std']} "
                                f"(범위 {b['min']}~{b['max']} 벗어남)")
                    rule[d] = b["std"]

    # ── max_count 교정 ──
    if not rule.get("max_count"):
        count_table = _MAX_COUNT_BY_CATEGORY.get(brand_category, _MAX_COUNT_CHARACTER_IP)
        base_count = count_table.get(obj_type, 1)
        rule["max_count"] = _scale_count(base_count, usable_area_m2)
        logger.info(f"[reference] {obj_type}: max_count 누락 → {rule['max_count']} "
                    f"(기준 {base_count}, 면적 {usable_area_m2:.0f}㎡)")



def _run_brand_agent(file_bytes: bytes, file_type: str = "pdf") -> dict:
    if file_type == "pptx":
        text = _extract_pptx_text(file_bytes)
    elif file_type == "docx":
        text = _extract_docx_text(file_bytes)
    elif file_type == "xlsx":
        text = _extract_xlsx_text(file_bytes)
    else:
        text = _extract_pdf_text(file_bytes)

    if not text.strip():
        return _fallback_brand_defaults()

    # Regex 주석 삽입 — LLM 단위 변환 실수 방지
    text = _annotate_measurements(text)

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        logger.info("[reference] API 키 없음 — 기본값 반환")
        return _fallback_brand_defaults()

    client = Anthropic(api_key=api_key)

    # LLM 설정은 app.llm_config 에서 중앙 관리 (temperature=0 강제 결정론)
    # 키 "small.reference" = 네임스페이스 규약. "small." prefix 제거 금지
    # (Shin이 "large.reference" 추가 시 충돌 방지). 상세: app/llm_config.py 최상단
    from app.llm_config import get_llm_config
    _cfg = get_llm_config("small.reference")

    # 2026-04-29: app.llm_harness.call_llm_tool_use 로 위임 — for-range retry + tool_use block
    # 추출 + token_tracker 수동 호출 → 하네스가 일괄. 응답 후처리 (defaults / clamp / normalize) 는 그대로.
    from app.nodes_small.llm_harness import call_llm_tool_use, LLMHarnessError

    try:
        result_obj, meta = call_llm_tool_use(
            client,
            model=_cfg["model"],
            max_tokens=_cfg["max_tokens"],
            temperature=_cfg["temperature"],
            system=[{
                "type": "text",
                "text": BRAND_SYSTEM,
                "cache_control": {"type": "ephemeral"},
            }],
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": BRAND_PROMPT, "cache_control": {"type": "ephemeral"}},
                    {"type": "text", "text": f"## 메뉴얼\n\n{text}"},
                ],
            }],
            tool_name="extract_brand_rules",
            tool_schema=BRAND_TOOL,
            response_model=BrandRulesResult,
            track_usage_node="reference",
            max_attempts=3,
        )
    except LLMHarnessError as e:
        # 2026-04-30: warning → error 격상. _MINIMAL_PLACEMENT_RULES_BY_CATEGORY 검증 도입 후
        # 3회 retry 모두 실패 = LLM 이 핵심 placement_rules 채우지 못함 = 배치 결과 품질 저하 확정.
        # silent 통과 차단 — 운영 모니터링에서 즉시 인지하고 매뉴얼/프롬프트 점검 트리거.
        logger.error(
            f"[reference] 하네스 실패 (3회 재시도 소진) — _fallback_brand_defaults() 사용. "
            f"placement_rules 부실 → 배치 품질 저하 가능. {type(e).__name__}: {e}"
        )
        return _fallback_brand_defaults()
    except Exception as e:
        logger.error(f"[reference] 예외 — _fallback_brand_defaults() 사용. 배치 품질 저하 가능: {e}")
        return _fallback_brand_defaults()

    # Pydantic → dict (기존 후처리 코드는 dict 접근). exclude_none=False, by_alias=False.
    raw = result_obj.model_dump()
    logger.info(f"[reference] brand 추출 완료 (attempts={meta.get('attempts', 1)})")

    # defaults merge
    for key, default in BRAND_DEFAULTS.items():
        if key in raw and isinstance(raw[key], dict):
            if raw[key].get("value") is None and default is not None:
                raw[key]["value"] = default
                raw[key]["source"] = "default"

    # placement_rules VMD_BOUNDARIES 기반 교정
    category = raw.get("brand_category", "기타")
    # 2026-05-01 SSOT trace: brand 추출 직후 카테고리 흐름 dump (라이브 검증용)
    from app.categories import dump_category_trace, get_category
    dump_category_trace(
        stage="reference.brand_extracted",
        raw_brand_category=category,
        placement_rules_count=len(raw.get("placement_rules", [])),
        placement_rule_types=[r.get("object_type") for r in raw.get("placement_rules", []) if isinstance(r, dict)],
        attempts=meta.get("attempts", 1),
    )

    # 2026-05-01 (G 옵션): placement_rules 부실 검사 — logger.warning 만, raise 안 함.
    # 부실해도 LLM 응답 (brand_category, clearspace_mm 등) 은 보존하고 다음 단계로 진행.
    # 이전 (PR #366) 의 model_validator raise 는 brand 전체 fallback 을 트리거하는 잘못된 결합이었음.
    # minimal_placement_rules 는 검증 트리거 X, 통계/모니터링 메타 (Category 주석 참조).
    _log_minimal_placement_rules_warning(category, raw.get("placement_rules", []))

    for rule in raw.get("placement_rules", []):
        _clamp_dimensions(rule, brand_category=category)

    # pair_rules: LLM 추출분 + VMD 기본 규칙 병합
    extracted_pairs = raw.get("pair_rules") or raw.get("relationships") or []
    merged_pairs = _merge_pair_rules(extracted_pairs)

    brand_fields = {k: v for k, v in raw.items()
                    if k not in ("placement_rules", "pair_rules", "relationships")}

    normalized_rules = _normalize_placement_rules(raw.get("placement_rules", []), brand_category=category)

    # character_bbox max_count 보정 (정규화 이후: 중복 합산 방지)
    # 1순위: max_count_source == "manual" → 매뉴얼 명시값 그대로 사용
    # 2순위: max_count_source == "inferred" + figures_mentioned 있음 → figures 개수로 교체
    # 3순위: max_count == null → _clamp_dimensions()가 면적 기준 기본값 적용
    figures = raw.get("figures_mentioned", [])
    for rule in normalized_rules:
        if rule.get("object_type") == "character_bbox":
            source = rule.get("max_count_source", "inferred")
            # b-3 후속: 자유 명명 (mooni_figure / stella_figure 등) 은 이미 별도 record 로 분리됨.
            # figures 기반 보정은 매뉴얼이 character_bbox 만 1개로 정의된 케이스에 한정 — label 이 std_id 와 동일할 때만.
            label = rule.get("label") or ""
            if source == "inferred" and figures and label in ("", "character_bbox"):
                rule["max_count"] = len(figures)
                logger.info(f"[reference] character_bbox max_count → {len(figures)} (figures_mentioned 기준)")

    return {
        "brand": brand_fields,
        "fire": FIRE_RULES,
        "construction": CONSTRUCTION_RULES,
        "placement_rules": normalized_rules,
        "pair_rules": merged_pairs,
    }


def _merge_pair_rules(extracted: list) -> list:
    """LLM 추출 pair_rules + VMD 기본 규칙 병합.

    LLM 추출분이 우선. 같은 (object_a, object_b) 쌍이 VMD에도 있으면 LLM 것을 유지.
    """
    # 추출분 구조화 (자연어 형식이면 변환)
    structured = []
    for rule in extracted:
        if isinstance(rule, dict) and "object_a" in rule and "relation" in rule:
            structured.append(rule)
        elif isinstance(rule, dict) and "rule" in rule:
            # Rendy 호환: {"rule": "자연어"} → 구조화 불가, 스킵
            logger.info(f"[reference] pair_rule 자연어 형식 스킵: {rule.get('rule', '')[:50]}")
        elif isinstance(rule, str):
            logger.info(f"[reference] pair_rule 문자열 형식 스킵: {rule[:50]}")

    # 추출분에 있는 쌍 키 수집
    extracted_pairs = set()
    for r in structured:
        extracted_pairs.add((r["object_a"], r["object_b"]))
        extracted_pairs.add((r["object_b"], r["object_a"]))  # 양방향

    # VMD 기본 규칙 중 추출분과 겹치지 않는 것만 추가
    for vmd_rule in VMD_PAIR_RULES:
        key = (vmd_rule["object_a"], vmd_rule["object_b"])
        if key not in extracted_pairs:
            structured.append({**vmd_rule, "source": "vmd_default"})

    logger.info(f"[reference] pair_rules: {len(extracted)}개 추출 → {len(structured)}개 (VMD 병합)")
    return structured


def _fallback_brand_defaults() -> dict:
    brand = {}
    for k, v in BRAND_DEFAULTS.items():
        brand[k] = {"value": v, "confidence": None, "source": "default"}
    brand["brand_category"] = "\uae30\ud0c0"
    # 2026-05-01 SSOT trace: fallback \uc9c4\uc785 \uba85\uc2dc dump (\ub77c\uc774\ube0c \uac80\uc99d \uc2dc fallback \ubc1c\ub3d9 \uc2dd\ubcc4)
    from app.categories import dump_category_trace
    dump_category_trace(
        stage="reference.fallback_used",
        raw_brand_category="\uae30\ud0c0",
        placement_rules_count=0,
        reason="LLM \ucd94\ucd9c \uc2e4\ud328 \ub610\ub294 brand_bytes \ubd80\uc7ac \u2014 _fallback_brand_defaults() \ud638\ucd9c",
    )
    return {
        "brand": brand,
        "fire": FIRE_RULES,
        "construction": CONSTRUCTION_RULES,
        "placement_rules": [],
        "pair_rules": [{**r, "source": "vmd_default"} for r in VMD_PAIR_RULES],
    }


def _normalize_placement_rules(rules: list, brand_category: str = "기타") -> list:
    """LLM이 추출한 placement_rules의 object_type을 표준 ID로 정규화.

    메뉴얼에 "포토존 배경 월"이라고 적혀있어도 → "photo_wall"으로 통일.
    "360도 포토존", "중앙 포토존" → "photo_island"로 분류.
    치수가 없으면 OBJECT_STANDARDS에서 기본값 채움.

    #472 b-3: raw 명명 (mooni_figure 등) 은 label 필드에 보존하면서 std_id 통합.
    """
    normalized = []
    seen_keys = set()  # (std_id, label) 복합 키 — 자유 명명 별도 record 보존
    for rule in rules:
        raw_type = rule.get("object_type", "")
        std_id = normalize_object_type(raw_type)
        std = OBJECT_STANDARDS.get(std_id)

        # b-3: raw 명명을 label 로 보존 (mooni_figure, stella_figure 등 자유 명명 매장 표시용).
        # name 필드 우선 (LLM 명시), fallback raw_type, 최후 std_id.
        rule["label"] = rule.get("name") or raw_type or std_id

        # 표준 치수 보충 + min/max 클램핑 — LLM 추출값이 범위 벗어나면 강제 교정
        rule["object_type"] = std_id
        bounds = get_vmd_boundaries(brand_category).get(std_id)
        if bounds:
            rule.setdefault("width_mm", bounds["width_mm"]["std"])
            rule.setdefault("depth_mm", bounds["depth_mm"]["std"])
            rule.setdefault("height_mm", bounds["height_mm"]["std"])
            rule.setdefault("front_edge", bounds.get("front_edge", "width"))
            # min/max 클램핑
            for dim in ("width_mm", "depth_mm", "height_mm"):
                val = rule.get(dim)
                b = bounds.get(dim)
                if val and b:
                    clamped = max(b["min"], min(val, b["max"]))
                    if clamped != val:
                        logger.info(f"[reference] {std_id} {dim}: {val} → {clamped} (min={b['min']}, max={b['max']})")
                        rule[dim] = clamped
        if std:
            rule.setdefault("name", std["name"])

        # wall_attachment: VMD 기본값으로 강제 — LLM 추출값이 flush 기물을 near로 바꾸는 것 방지
        from app.vmd_constants import VMD_WALL_ATTACHMENT
        vmd_wa = VMD_WALL_ATTACHMENT.get(std_id)
        if vmd_wa and rule.get("wall_attachment") != vmd_wa:
            logger.info(f"[reference] {std_id} wall_attachment: {rule.get('wall_attachment')} → {vmd_wa} (VMD 강제)")
            rule["wall_attachment"] = vmd_wa

        # front_edge: VMD 기본값으로 강제 — LLM이 임의로 바꾸는 것 방지
        if bounds and bounds.get("front_edge"):
            vmd_fe = bounds["front_edge"]
            if rule.get("front_edge") != vmd_fe:
                logger.info(f"[reference] {std_id} front_edge: {rule.get('front_edge')} → {vmd_fe} (VMD 강제)")
                rule["front_edge"] = vmd_fe

        # b-3 + e: 복합 키 (std_id, label) 로 중복 판정.
        # - 같은 std_id + 같은 label → 변형 명명 동일 의미 (consultation_desk vs 한글 "상담 데스크"),
        #   max_count 합산 X, max 적용 (의미 동일 → 더 큰 수가 정답)
        # - 같은 std_id + 다른 label → 자유 명명 다른 개체 (mooni_figure vs stella_figure),
        #   별도 record 유지하여 6종 캐릭터 모두 매장에 보존
        key = (std_id, rule["label"])
        if key in seen_keys:
            for existing in normalized:
                if existing["object_type"] == std_id and existing.get("label") == rule["label"]:
                    existing["max_count"] = max(
                        existing.get("max_count", 1),
                        rule.get("max_count", 1),
                    )
                    break
        else:
            seen_keys.add(key)
            normalized.append(rule)

    logger.info(f"[reference] placement_rules 정규화: {len(rules)}개 → {len(normalized)}개 (std_id, label) 키")
    return normalized


def _annotate_measurements(text: str) -> str:
    """PDF 텍스트에 단위 변환 주석 삽입 (Regex 단계).

    LLM이 단위 변환 실수를 줄이도록 원본 표현 옆에 mm 변환값을 추가.
    순서 중요: 범위 패턴 먼저, 단일 패턴 나중에 처리.

    예: "50cm"      → "50cm[=500mm]"
        "1.2m"      → "1.2m[=1200mm]"
        "30~50cm"   → "30~50cm[=300~500mm, 최솟값=300mm]"
        "300~500mm" → "300~500mm[최솟값=300mm]"
    """
    # 1. cm 범위: "30~50cm" → "30~50cm[=300~500mm, 최솟값=300mm]"
    def replace_cm_range(m):
        lo_mm = int(float(m.group(1)) * 10)
        hi_mm = int(float(m.group(2)) * 10)
        return f"{m.group(0)}[={lo_mm}~{hi_mm}mm, 최솟값={lo_mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*[~～]\s*(\d+(?:\.\d+)?)\s*cm(?!\[)', replace_cm_range, text)

    # 2. m 범위 (mm 제외): "1~2m" → "1~2m[=1000~2000mm, 최솟값=1000mm]"
    def replace_m_range(m):
        lo_mm = int(float(m.group(1)) * 1000)
        hi_mm = int(float(m.group(2)) * 1000)
        return f"{m.group(0)}[={lo_mm}~{hi_mm}mm, 최솟값={lo_mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*[~～]\s*(\d+(?:\.\d+)?)\s*m(?!m|\[|²)', replace_m_range, text)

    # 3. mm 범위: "300~500mm" → "300~500mm[최솟값=300mm]"
    def replace_mm_range(m):
        lo_mm = int(m.group(1))
        return f"{m.group(0)}[최솟값={lo_mm}mm]"
    text = re.sub(r'(\d+)\s*[~～]\s*(\d+)\s*mm(?!\[)', replace_mm_range, text)

    # 4. 단일 cm: "50cm" → "50cm[=500mm]"
    def replace_cm(m):
        mm = int(float(m.group(1)) * 10)
        return f"{m.group(0)}[={mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*cm(?!\[)', replace_cm, text)

    # 5. 단일 m (mm·m²·주석 제외): "1.2m" → "1.2m[=1200mm]"
    def replace_m(m):
        mm = int(float(m.group(1)) * 1000)
        return f"{m.group(0)}[={mm}mm]"
    text = re.sub(r'(\d+(?:\.\d+)?)\s*m(?!m|\[|²)', replace_m, text)

    return text


def _extract_pdf_text(pdf_bytes: bytes) -> str:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page_texts = []
    for page in doc:
        texts = []
        table_bboxes = []

        try:
            for table in page.find_tables():
                table_bboxes.append(table.bbox)
                rows = []
                for row in table.extract():
                    cells = [str(c).strip() for c in row if c and str(c).strip()]
                    if cells:
                        rows.append(" ".join(cells))
                if rows:
                    texts.append("\n".join(rows))
        except Exception:
            pass

        for block in page.get_text("blocks", sort=True):
            if block[6] != 0:
                continue
            text = block[4].strip()
            if not text:
                continue
            bx0, by0, bx1, by1 = block[:4]
            in_table = any(
                bx0 < tx1 and bx1 > tx0 and by0 < ty1 and by1 > ty0
                for (tx0, ty0, tx1, ty1) in table_bboxes
            )
            if not in_table:
                texts.append(text)

        if texts:
            page_texts.append("\n".join(texts))

    doc.close()
    return "\n\n".join(page_texts)


def _extract_pptx_text(data: bytes) -> str:
    import io
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _extract_shapes(shapes, texts):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _extract_shapes(shape.shapes, texts)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line:
                        texts.append(line)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" ".join(cells))

    prs = Presentation(io.BytesIO(data))
    texts = []
    for slide in prs.slides:
        _extract_shapes(slide.shapes, texts)
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            note = notes_tf.text.strip()
            if note:
                texts.append(note)
        except Exception:
            pass
    return "\n".join(texts)


def _extract_docx_text(data: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(data))
    texts = []
    for section in doc.sections:
        for hf in (section.header, section.footer):
            for para in hf.paragraphs:
                if para.text.strip():
                    texts.append(para.text.strip())
    texts += [para.text for para in doc.paragraphs if para.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                texts.append(" ".join(cells))
    return "\n".join(texts)


def _extract_xlsx_text(data: bytes) -> str:
    import io
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    texts = []
    for ws in wb.worksheets:
        texts.append(f"[{ws.title}]")
        for row in ws.iter_rows(values_only=True):
            line = " ".join(str(cell) for cell in row if cell is not None and str(cell).strip())
            if line:
                texts.append(line)
    wb.close()
    return "\n".join(texts)
